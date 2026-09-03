#!/usr/bin/env python3
"""Deterministic static PPTX house-rule checks."""
from __future__ import annotations

import argparse
import ast
import hashlib
import io
import json
import math
import posixpath
import re
import sys
import zipfile
from dataclasses import asdict, dataclass
from decimal import Decimal, InvalidOperation
from pathlib import Path
from xml.etree import ElementTree as ET

import yaml
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from openpyxl import load_workbook

if sys.version_info < (3, 11):
    raise SystemExit("audit.py requires Python 3.11 or newer (see DEVELOPMENT_PLAN.md section 4)")

from schemas.manifest import validate as validate_manifest

@dataclass(frozen=True)
class Issue:
    rule: str
    slide: int
    shape: str
    evidence: str


def load_rules(path):
    with path.open(encoding="utf-8") as handle:
        rules = yaml.safe_load(handle)
    required = {"units", "notation", "numeric_tokens", "manifest", "issues", "styles", "default_style"}
    missing = sorted(required - rules.keys())
    if missing:
        raise ValueError(f"house-rules missing sections: {', '.join(missing)}")
    if not isinstance(rules["styles"], dict) or not rules["styles"]:
        raise ValueError("house-rules styles must contain at least one style")
    return rules


def style_rules(rules, manifest_path=None):
    """Return the rules for a manifest's style, or the configured default.

    A missing/unknown style is an error: silently falling back to another
    style would produce a PASS against the wrong house standard.
    """
    style = None
    if manifest_path is not None:
        payload = json.loads(manifest_path.read_text(encoding="utf-8"))
        style = payload.get("style")
        if not style:
            raise ValueError("manifest style is required")
    else:
        raise ValueError("style is unknown: manifest is required")
    styles = rules.get("styles", {})
    if style not in styles:
        raise ValueError(f"unknown style {style!r}; expected one of {sorted(styles)}")
    required_style = {
        "fonts", "sizes", "table", "zones", "forbidden", "palette", "palette_usage",
        "components", "qa", "role_min_pt", "layout",
    }
    missing = sorted(required_style - styles[style].keys())
    if missing:
        raise ValueError(f"style {style!r} missing sections: {', '.join(missing)}")
    def merge(base, override):
        merged = dict(base)
        for key, value in override.items():
            if isinstance(value, dict) and isinstance(merged.get(key), dict):
                merged[key] = merge(merged[key], value)
            else:
                merged[key] = value
        return merged

    effective = merge(rules, styles[style])
    return effective


def check_preflight_alignment(rules):
    """Detect drift in the overlapping checks owned by the skill preflight."""
    source_path = Path(__file__).with_name("skill") / "shin-ppt1" / "scripts" / "preflight.py"
    if not source_path.is_file():
        return []
    tree = ast.parse(source_path.read_text(encoding="utf-8"), filename=str(source_path))
    constants = {}
    for node in tree.body:
        if isinstance(node, ast.Assign):
            for target in node.targets:
                if isinstance(target, ast.Name) and target.id in {"FONTS_OK"}:
                    constants[target.id] = ast.literal_eval(node.value)
    issues = []
    expected_fonts = {rules["fonts"]["body"], rules["fonts"]["heading"]}
    if constants.get("FONTS_OK") is not None and set(constants["FONTS_OK"]) != expected_fonts:
        issues.append(Issue(
            "contract.preflight_fonts", 0, "-",
            f"preflight={sorted(constants['FONTS_OK'])}, house-rules={sorted(expected_fonts)}",
        ))
    tolerance = rules["qa"].get("colw_tolerance_emu")
    match = re.search(r"abs\(total\s*-\s*cx\)\s*>\s*(\d+)", source_path.read_text(encoding="utf-8"))
    if tolerance is not None and match and float(tolerance) != float(match.group(1)):
        issues.append(Issue(
            "contract.preflight_colw_tolerance", 0, "-",
            f"preflight={match.group(1)}EMU, house-rules={tolerance}EMU",
        ))
    return issues


def text_shapes(slide):
    return (shape for shape in slide.shapes if getattr(shape, "has_text_frame", False))


def text(shape):
    return "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs)


def fonts(text_frame):
    return {
        run.font.name
        for paragraph in text_frame.paragraphs
        for run in paragraph.runs
        if run.font.name
    }


def check_fonts(prs, rules):
    allowed = {rules["fonts"]["heading"], rules["fonts"]["body"]}
    issues = []
    for page, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            frames = []
            if getattr(shape, "has_text_frame", False):
                frames.append(shape.text_frame)
            if getattr(shape, "has_table", False):
                frames.extend(cell.text_frame for row in shape.table.rows for cell in row.cells)
            extra = sorted(set().union(*(fonts(frame) for frame in frames)) - allowed) if frames else []
            if extra:
                issues.append(Issue("forbidden.third_font", page, shape.name,
                                    f"허용되지 않은 글꼴: {', '.join(extra)}"))
    return issues


def check_notation(prs, rules):
    forbidden = tuple(rules["notation"]["negative_forbidden"])
    issues = []
    for page, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            values = []
            if getattr(shape, "has_text_frame", False):
                values.append(text(shape))
            if getattr(shape, "has_table", False):
                values.extend(cell.text for row in shape.table.rows for cell in row.cells)
            found = sorted({mark for value in values for mark in forbidden if mark in value})
            if found:
                issues.append(Issue("notation.negative_forbidden", page, shape.name,
                                    f"금지 음수 부호: {', '.join(found)}"))
    return issues


def run_color(run):
    try:
        rgb = run.font.color.rgb
        return str(rgb).upper() if rgb is not None else None
    except (AttributeError, TypeError, ValueError):
        return None


def check_negative_red(prs, rules):
    if "negative_red" not in rules["forbidden"]:
        return []
    negative = re.escape(rules["notation"]["negative"])
    pattern = re.compile(rf"{negative}\s*\d")
    red = str(rules["palette"]["red"]).upper()
    issues = []
    for page, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            frames = []
            if getattr(shape, "has_text_frame", False):
                frames.append(shape.text_frame)
            if getattr(shape, "has_table", False):
                frames.extend(cell.text_frame for row in shape.table.rows for cell in row.cells)
            for frame in frames:
                for paragraph in frame.paragraphs:
                    for run in paragraph.runs:
                        if pattern.search(run.text) and run_color(run) == red:
                            issues.append(Issue("forbidden.negative_red", page, shape.name,
                                                f"빨강 음수: {run.text!r}"))
    return issues


def check_red_runs_per_line(prs, rules):
    section = "palette_" + "usage"
    maximum = int(rules[section]["red_max_per_line"])
    red = str(rules["palette"]["red"]).upper()
    issues = []
    for page, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            frames = []
            if getattr(shape, "has_text_frame", False):
                frames.append(shape.text_frame)
            if getattr(shape, "has_table", False):
                frames.extend(cell.text_frame for row in shape.table.rows for cell in row.cells)
            for frame in frames:
                for line, paragraph in enumerate(frame.paragraphs, 1):
                    count = sum(1 for run in paragraph.runs
                                if run.text.strip() and run_color(run) == red)
                    if count > maximum:
                        issues.append(Issue(
                            section + ".red_max_per_line", page, shape.name,
                            f"{line}번째 문단 빨강 런={count} > {maximum}",
                        ))
    return issues


def minimum_font_size(shape_name, rules, is_table=False):
    sizes = rules["sizes"]
    if is_table:
        return float(sizes["table_body_min_pt"])
    base = shape_name.split("#", 1)[0]
    table = rules["role_min_pt"]
    key = table.get(base)
    if key is None:
        if "/" in base:
            raise ValueError(f"role_min_pt[{base!r}] is not defined for this style")
        key = table["_claim_shape"]
    if key not in sizes:
        raise ValueError(f"role_min_pt[{base!r}] references missing sizes key: {key!r}")
    return float(sizes[key])


def check_font_sizes(prs, rules):
    issues = []
    for page, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            frames = []
            if getattr(shape, "has_text_frame", False):
                frames.append((shape.text_frame, False))
            if getattr(shape, "has_table", False):
                frames.extend((cell.text_frame, True) for row in shape.table.rows for cell in row.cells)
            for frame, is_table in frames:
                if not any(run.text.strip() for paragraph in frame.paragraphs for run in paragraph.runs):
                    continue
                minimum = minimum_font_size(shape.name, rules, is_table)
                for paragraph in frame.paragraphs:
                    for run in paragraph.runs:
                        if not run.text.strip() or run.font.size is None:
                            continue
                        actual = run.font.size.pt
                        if actual + 0.001 < minimum:
                            issues.append(Issue("sizes.body_min_pt", page, shape.name,
                                                f"{actual:g}pt < 역할별 하한 {minimum:g}pt: {run.text[:30]!r}"))
    return issues


ALIGN_NAMES = {
    None: "left",
    PP_ALIGN.LEFT: "left",
    PP_ALIGN.CENTER: "center",
    PP_ALIGN.RIGHT: "right",
    PP_ALIGN.JUSTIFY: "justify",
}


def numeric_only(value, rules):
    candidate = value.strip()
    for marker in rules["notation"]["negative_forbidden"]:
        if candidate.startswith(marker):
            candidate = rules["notation"]["negative"] + candidate[len(marker):].lstrip()
            break
    if not re.fullmatch(rules["numeric_tokens"]["pattern"], candidate):
        return False
    number = candidate
    for sign in (rules["notation"]["negative"], rules["notation"]["positive"]):
        if number.startswith(sign):
            number = number[len(sign):]
            break
    if number.endswith("%"):
        number = number[:-1]
    number = number.replace(rules["notation"]["thousands_sep"], "")
    number = number.replace(rules["notation"]["decimal_sep"], ".")
    try:
        Decimal(number)
    except InvalidOperation:
        return False
    return True


def is_subtotal_row(row, rules):
    if rules["table"].get("subtotal_row_style") != "bold_no_fill":
        return False
    nonempty = [cell for cell in row.cells if cell.text.strip()]
    return bool(nonempty) and all(
        cell.fill.type is None
        and all(run.font.bold is True
                for paragraph in cell.text_frame.paragraphs
                for run in paragraph.runs if run.text.strip())
        for cell in nonempty
    )


def check_table_alignments(prs, rules):
    table_rules = rules["table"]
    header_align = table_rules["header_align"]
    default_align = table_rules["default_align"]
    long_text_align = table_rules["long_text_col_align"]
    numeric_align = table_rules["numeric_col_align"]
    issues = []
    for page, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False) or not shape.table.rows:
                continue
            body_rows = list(shape.table.rows)[1:]
            long_text_columns = {
                col_index
                for col_index in range(len(shape.table.columns))
                if (values := [row.cells[col_index].text.strip() for row in body_rows
                               if row.cells[col_index].text.strip()])
                and all(not numeric_only(value, rules) for value in values)
            }
            for row_index, row in enumerate(shape.table.rows):
                subtotal = row_index > 0 and is_subtotal_row(row, rules)
                for col_index, cell in enumerate(row.cells):
                    value = cell.text.strip()
                    if not value:
                        continue
                    if row_index == 0:
                        rule, allowed = "table.header_align", {header_align}
                    elif numeric_only(value, rules):
                        rule = "table.numeric_col_align"
                        allowed = {numeric_align, default_align} if subtotal else {numeric_align}
                    elif col_index in long_text_columns:
                        rule, allowed = "table.default_align", {default_align, long_text_align}
                    else:
                        rule, allowed = "table.default_align", {default_align}
                    actual = {
                        ALIGN_NAMES.get(paragraph.alignment, str(paragraph.alignment).lower())
                        for paragraph in cell.text_frame.paragraphs if paragraph.text.strip()
                    }
                    bad = sorted(actual - allowed)
                    if bad:
                        issues.append(Issue(
                            rule, page, shape.name,
                            f"셀[{row_index},{col_index}]={value!r}: align={bad}, expected={sorted(allowed)}",
                        ))
    return issues


def check_footnotes(prs, rules):
    zones = rules["zones"]
    markers = tuple(rules["components"]["footnote_markers"])
    tolerance = rules["qa"]["canvas_overflow_tolerance_in"]
    issues = []
    for page, slide in enumerate(prs.slides, 1):
        for shape in text_shapes(slide):
            if not shape.name.startswith("footer/"):
                continue
            paragraphs = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
            if not paragraphs or not paragraphs[0].startswith(markers):
                continue
            expected = zones["footnote_bottom_y"] - zones["footnote_line_step"] * len(paragraphs)
            actual = shape.top / rules["units"]["emu_per_inch"]
            if abs(actual - expected) > tolerance:
                issues.append(Issue("zones.footnote_bottom_y", page, shape.name,
                                    f"각주 y={actual:.2f}, 기대값={expected:.2f}"))
    return issues


def estimated_height(shape, rules):
    width_pt = shape.width / rules["units"]["emu_per_inch"] * rules["units"]["pt_per_inch"]
    if width_pt <= 0:
        return math.inf
    sizes = rules["sizes"]
    baseline = sizes.get("body_min_pt", sizes.get("card_body_pt", sizes.get("table_body_min_pt")))
    if baseline is None:
        raise ValueError("style sizes must define a text-size baseline")
    lines, max_size = 0, float(baseline)
    for paragraph in shape.text_frame.paragraphs:
        weighted = 0.0
        for run in paragraph.runs:
            size = run.font.size.pt if run.font.size else max_size
            max_size = max(max_size, size)
            weighted += sum(size * (1.0 if ord(char) > 127 else 0.55) for char in run.text)
        lines += max(1, math.ceil(weighted / width_pt))
    return lines * max_size * 1.15 / rules["units"]["pt_per_inch"]


def check_overflow(prs, rules):
    tolerance = rules["qa"]["canvas_overflow_tolerance_in"]
    issues = []
    for page, slide in enumerate(prs.slides, 1):
        for shape in text_shapes(slide):
            if not text(shape).strip():
                continue
            needed = estimated_height(shape, rules)
            available = shape.height / rules["units"]["emu_per_inch"]
            if needed > available + tolerance:
                issues.append(Issue("qa.text_max_ymax_pt", page, shape.name,
                                    f"정적 근사 필요높이={needed:.2f}in, 상자높이={available:.2f}in"))
    return issues


def check_title_right(prs, rules):
    if not rules["zones"].get("title_right_clear", False):
        return []
    title, layout = rules["components"]["page_title"], rules["layout"]
    right_start = layout["width"] - 3.1
    y0, y1 = title["y"], title["y"] + title["h"]
    issues = []
    for page, slide in enumerate(prs.slides, 1):
        for shape in text_shapes(slide):
            unit = rules["units"]["emu_per_inch"]
            x, y, h = shape.left / unit, shape.top / unit, shape.height / unit
            canonical_title = abs(x - layout["margin_x"]) < 0.01 and abs(y - y0) < 0.01
            if not canonical_title and x >= right_start and y < y1 and y + h > y0:
                issues.append(Issue("zones.title_right_clear", page, shape.name,
                                    f"제목 우상단 텍스트: {text(shape)[:60]}"))
    return issues


def check_table_geometry(prs, rules):
    unit = rules["units"]["emu_per_inch"]
    tolerance = rules["qa"]["colw_tolerance_emu"]
    issues = []
    for page, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            table = shape.table
            if rules["table"].get("colw_sum_must_equal_width", False):
                total = sum(column.width for column in table.columns)
                if abs(total - shape.width) > tolerance:
                    issues.append(Issue("table.colw_sum_must_equal_width", page, shape.name,
                                        f"열 합={total}EMU, 표 폭={shape.width}EMU"))
            for row_index, row in enumerate(table.rows, 1):
                two_line = any("\n" in cell.text or "\r" in cell.text for cell in row.cells)
                key = "row_height_2line_min" if two_line else "row_height_min"
                if key not in rules["table"]:
                    continue
                minimum = float(rules["table"][key])
                actual = row.height / unit
                if actual + rules["units"]["epsilon_in"] < minimum:
                    issues.append(Issue(f"table.{key}", page, shape.name,
                                        f"{row_index}행 높이={actual:.2f}in < {minimum:.2f}in"))
    return issues


def check_canvas_and_content(prs, rules):
    unit = rules["units"]["emu_per_inch"]
    tolerance = rules["qa"]["canvas_overflow_tolerance_in"]
    page_width, page_height = rules["layout"]["width"], rules["layout"]["height"]
    content_max = rules["zones"]["content_max_y"]
    content_exempt = {"header/draft_tag"}
    issues = []
    for page, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            x, y = shape.left / unit, shape.top / unit
            right, bottom = x + shape.width / unit, y + shape.height / unit
            if x < -tolerance or y < -tolerance or right > page_width + tolerance or bottom > page_height + tolerance:
                issues.append(Issue("layout.canvas_overflow", page, shape.name,
                                    f"bounds=({x:.2f},{y:.2f},{right:.2f},{bottom:.2f})in"))
            exempt = shape.name in content_exempt or shape.name.startswith("footer/")
            if not exempt and bottom > content_max + tolerance:
                issues.append(Issue("zones.content_max_y", page, shape.name,
                                    f"본문 하단={bottom:.2f}in > {content_max:.2f}in"))
    return issues


def shape_role(shape):
    return shape.name.split("#", 1)[0]


def check_chip_geometry(prs, rules):
    unit = rules["units"]["emu_per_inch"]
    tolerance = 10 ** -int(rules["units"]["bounds_round_in"])
    issues = []
    for page, slide in enumerate(prs.slides, 1):
        backgrounds = [shape for shape in slide.shapes if shape_role(shape) == "chip/bg"]
        descriptions = [shape for shape in slide.shapes if shape_role(shape) == "chip/desc"]
        if not backgrounds and not descriptions:
            continue
        chip_rules = rules["components"].get("chip")
        if chip_rules is None:
            raise ValueError("components.chip rules are required when chip shapes exist")
        desc_gap = float(chip_rules["desc_gap"])
        desc_max_width = float(chip_rules["desc_w"])
        right_limit = float(rules["layout"]["width"]) - float(rules["layout"]["margin_x"])
        unused = set(range(len(backgrounds)))
        for desc in sorted(descriptions, key=lambda shape: (shape.top, shape.left)):
            dx = desc.left / unit
            dy0, dy1 = desc.top / unit, (desc.top + desc.height) / unit
            candidates = []
            for index in unused:
                bg = backgrounds[index]
                bg_right = (bg.left + bg.width) / unit
                by0, by1 = bg.top / unit, (bg.top + bg.height) / unit
                overlaps_y = dy0 < by1 + tolerance and dy1 > by0 - tolerance
                if overlaps_y and dx >= bg_right - tolerance:
                    candidates.append((abs(dx - bg_right - desc_gap), index, bg_right))
            if not candidates:
                issues.append(Issue(
                    "layout.chip_pair", page, desc.name,
                    "세로로 겹치면서 왼쪽에 있는 chip/bg를 찾지 못함",
                ))
                continue
            _, index, bg_right = min(candidates)
            unused.remove(index)
            actual_gap = dx - bg_right
            width = desc.width / unit
            canvas_width = right_limit - dx
            if abs(actual_gap - desc_gap) > tolerance:
                issues.append(Issue(
                    "layout.chip_desc_gap", page, desc.name,
                    f"gap={actual_gap:.4f}in, expected={desc_gap:.4f}in",
                ))
            if width > desc_max_width + tolerance:
                issues.append(Issue(
                    "layout.chip_desc_width", page, desc.name,
                    f"w={width:.4f}in > desc_w={desc_max_width:.4f}in",
                ))
            if width > canvas_width + tolerance:
                issues.append(Issue(
                    "layout.chip_desc_canvas", page, desc.name,
                    f"w={width:.4f}in > right_limit-desc.x={canvas_width:.4f}in",
                ))
    return issues


def format_number(value, display, rules):
    rounding = display.get("rounding")
    if rounding is None:
        return str(value)
    number = float(value)
    absolute = f"{abs(number):,.{rounding}f}"
    absolute = absolute.replace(",", "\0").replace(".", rules["notation"]["decimal_sep"])
    absolute = absolute.replace("\0", rules["notation"]["thousands_sep"])
    if number < 0:
        return rules["notation"]["negative"] + absolute
    if str(display.get("text", "")).startswith(rules["notation"].get("positive", "+")):
        return rules["notation"].get("positive", "+") + absolute
    return absolute


def cell_value(sheet, ref):
    return sheet[ref].value


def calculate_claim(claim, workbook):
    source, transform = claim["source"], claim["transform"]
    sheet = workbook[source["sheet"]]
    kind = transform["type"]
    if kind == "identity":
        return cell_value(sheet, source["ref"])
    if kind == "sum":
        return sum(cell.value or 0 for row in sheet[transform["range"]] for cell in row)
    if kind == "ratio":
        return cell_value(sheet, transform["numerator"]) / cell_value(sheet, transform["denominator"])
    if kind == "delta":
        return cell_value(sheet, transform["to"]) - cell_value(sheet, transform["from"])
    if kind == "cagr":
        start, end = cell_value(sheet, transform["start"]), cell_value(sheet, transform["end"])
        return (end / start) ** (1 / transform["periods"]) - 1
    if kind == "unverified":
        return None
    raise ValueError(f"unsupported transform: {kind}")


CHART_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
OFFICE_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def normalized_number(value):
    """Return a stable numeric identity without conflating bools with 0/1."""
    if isinstance(value, bool) or value is None:
        return None
    if isinstance(value, (int, float, Decimal)):
        try:
            number = Decimal(str(value))
        except InvalidOperation:
            return None
        return number.normalize() if number.is_finite() else None
    return None


def related_part_name(part_name, target):
    if target.startswith("/"):
        return target.lstrip("/")
    return posixpath.normpath(posixpath.join(posixpath.dirname(part_name), target))


def relationships(archive, part_name):
    directory, filename = posixpath.split(part_name)
    rels_name = posixpath.join(directory, "_rels", filename + ".rels")
    if rels_name not in archive.namelist():
        return []
    root = ET.fromstring(archive.read(rels_name))
    return [
        {
            "id": node.attrib["Id"],
            "type": node.attrib["Type"],
            "target": related_part_name(part_name, node.attrib["Target"]),
        }
        for node in root.findall(f"{{{REL_NS}}}Relationship")
        if node.attrib.get("TargetMode") != "External"
    ]


def chart_series_references(pptx_path):
    """Yield slide/chart/series value ranges from embedded chart workbooks.

    Only c:val, c:yVal, and c:bubbleSize are data series.  c:cat and c:xVal
    are deliberately excluded because they contain category/axis labels.
    """
    with zipfile.ZipFile(pptx_path) as archive:
        slide_names = sorted(
            (name for name in archive.namelist()
             if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)),
            key=lambda name: int(re.search(r"slide(\d+)\.xml$", name).group(1)),
        )
        for slide_name in slide_names:
            page = int(re.search(r"slide(\d+)\.xml$", slide_name).group(1))
            chart_rels = [rel for rel in relationships(archive, slide_name)
                          if rel["type"].endswith("/chart")]
            for slide_rel in chart_rels:
                chart_name = slide_rel["target"]
                chart_root = ET.fromstring(archive.read(chart_name))
                package_rels = [rel for rel in relationships(archive, chart_name)
                                if rel["type"].endswith("/package")]
                external = chart_root.find(f".//{{{CHART_NS}}}externalData")
                rel_id = external.attrib.get(f"{{{OFFICE_REL_NS}}}id") if external is not None else None
                workbook_rel = next((rel for rel in package_rels if rel["id"] == rel_id), None)
                if workbook_rel is None and len(package_rels) == 1:
                    workbook_rel = package_rels[0]
                if workbook_rel is None:
                    raise ValueError(f"chart workbook relationship missing: {chart_name}")
                workbook = load_workbook(
                    io.BytesIO(archive.read(workbook_rel["target"])),
                    data_only=True,
                    read_only=True,
                )
                try:
                    for series_index, series in enumerate(
                            chart_root.findall(f".//{{{CHART_NS}}}ser"), 1):
                        for value_tag in ("val", "yVal", "bubbleSize"):
                            formula = series.find(
                                f"{{{CHART_NS}}}{value_tag}/{{{CHART_NS}}}numRef/{{{CHART_NS}}}f"
                            )
                            if formula is None or not formula.text:
                                continue
                            match = re.fullmatch(r"(?:'((?:[^']|'')+)'|([^!]+))!(.+)", formula.text)
                            if not match:
                                raise ValueError(
                                    f"unsupported chart series formula {formula.text!r}: {chart_name}"
                                )
                            sheet_name = (match.group(1) or match.group(2)).replace("''", "'")
                            cell_range = match.group(3).replace("$", "")
                            if sheet_name not in workbook.sheetnames:
                                raise ValueError(
                                    f"chart worksheet missing {sheet_name!r}: {chart_name}"
                                )
                            cells = workbook[sheet_name][cell_range]
                            if not isinstance(cells, tuple):
                                cells = ((cells,),)
                            elif cells and not isinstance(cells[0], tuple):
                                cells = (cells,)
                            values = [cell.value for row in cells for cell in row
                                      if normalized_number(cell.value) is not None]
                            yield page, chart_name, series_index, formula.text, values
                finally:
                    workbook.close()


def claim_source_numbers(claims, workbooks, resolved_sources):
    """Index numeric evidence by slide from source sheets named by claims."""
    indexed = {}
    allowed_by_page = {}
    for claim in claims:
        source = claim.get("source", {})
        source_path = resolved_sources.get(id(claim))
        sheet_name = source.get("sheet")
        if not source_path or source_path not in workbooks or sheet_name not in workbooks[source_path].sheetnames:
            continue
        key = (source_path, sheet_name)
        if key not in indexed:
            indexed[key] = {
                normalized_number(cell.value)
                for row in workbooks[source_path][sheet_name].iter_rows()
                for cell in row
                if normalized_number(cell.value) is not None
            }
        for page in {int(item["slide"]) for item in claim.get("placements", [])}:
            allowed_by_page.setdefault(page, set()).update(indexed[key])
    return allowed_by_page


def check_chart_series(pptx_path, claims, workbooks, resolved_sources):
    allowed_by_page = claim_source_numbers(claims, workbooks, resolved_sources)
    issues = []
    for page, chart_name, series_index, formula, values in chart_series_references(pptx_path):
        allowed = allowed_by_page.get(page, set())
        for point_index, value in enumerate(values, 1):
            if normalized_number(value) not in allowed:
                issues.append(Issue(
                    "claim.unregistered_chart_series_value", page, chart_name,
                    f"계열 {series_index} 값 {point_index} ({formula})={value!r}: "
                    "manifest가 가리키는 원천 시트에 없음",
                ))
    return issues


def check_claims(prs, rules, manifest_path, source_root=None, pptx_path=None):
    payload = json.loads(manifest_path.read_text(encoding="utf-8"))
    schema_errors = validate_manifest(payload, rules)
    if schema_errors:
        raise ValueError("manifest schema: " + " | ".join(schema_errors))
    claims = payload["claims"]
    issues, changes, warnings = [], [], []
    workbooks = {}
    resolved_sources = {}
    root = source_root or manifest_path.parent
    for claim in claims:
        shape_id = claim["shape_id"]
        display = str(claim["display"]["text"])
        transform = claim["transform"]["type"]
        placements = claim["placements"]
        for placement in placements:
            page = int(placement["slide"])
            placed_text = str(placement["text"])
            if display not in placed_text:
                issues.append(Issue("claim.cross_page_consistency", page, shape_id,
                                    f"placement={placed_text!r}에 display={display!r} 없음"))
            if page < 1 or page > len(prs.slides):
                issues.append(Issue("claim.source_manifest_pptx", page, shape_id,
                                    "placement 슬라이드 없음"))
                continue
            slide = prs.slides[page - 1]
            if placement["type"] == "shape":
                matches = [shape for shape in slide.shapes if shape.name == placement["name"]]
                if len(matches) != 1:
                    issues.append(Issue("claim.source_manifest_pptx", page, shape_id,
                                        f"named shape 개수={len(matches)}, 기대값=1"))
                    continue
                shape = matches[0]
                actual_text = text(shape).strip()
                if actual_text != placed_text:
                    issues.append(Issue("claim.source_manifest_pptx", page, shape_id,
                                        f"placement={placed_text!r}, pptx={actual_text!r}"))
                check_shape_placement(shape, placement, rules, page, shape_id, issues)
            elif placement["type"] == "cell":
                tables = [shape for shape in slide.shapes
                          if getattr(shape, "has_table", False) and shape.name == placement["table"]]
                if len(tables) != 1:
                    issues.append(Issue("claim.source_manifest_pptx", page, shape_id,
                                        f"named table 개수={len(tables)}, 기대값=1"))
                    continue
                try:
                    actual_text = tables[0].table.cell(int(placement["row"]), int(placement["col"])).text
                except IndexError:
                    issues.append(Issue("claim.source_manifest_pptx", page, shape_id, "표 셀 좌표 범위 초과"))
                    continue
                if actual_text != placed_text:
                    issues.append(Issue("claim.source_manifest_pptx", page, shape_id,
                                        f"placement={placed_text!r}, pptx cell={actual_text!r}"))
            else:
                raise ValueError(f"claim[{shape_id}] unknown placement type: {placement['type']}")

        source = claim.get("source", {})
        if transform == "unverified":
            warnings.append(Issue(
                "calc.unverified_claim", int(claim.get("slide", placements[0]["slide"])), shape_id,
                f"value={display!r}, note={claim['transform']['note']!r}: 원천 계산 검증 안 됨",
            ))
            continue
        source_path = root / source["file"]
        if not source_path.exists():
            raise FileNotFoundError(f"claim[{shape_id}] source missing: {source_path}")
        digest = "sha256:" + hashlib.sha256(source_path.read_bytes()).hexdigest()
        if source.get("file_hash") != digest:
            issues.append(Issue("claim.source_manifest_pptx", page, shape_id,
                                "원천 파일 해시 불일치"))
        if source_path not in workbooks:
            workbooks[source_path] = load_workbook(source_path, data_only=True, read_only=True)
        resolved_sources[id(claim)] = source_path
        calculated = calculate_claim(claim, workbooks[source_path])
        expected = format_number(calculated, claim["display"], rules)
        override = claim.get("override")
        if override:
            source_display = expected
            if claim["display"].get("rounding") is None and isinstance(calculated, (int, float)):
                decimal_sep = rules["notation"]["decimal_sep"]
                decimals = len(display.rsplit(decimal_sep, 1)[1]) if decimal_sep in display else 0
                source_display = format_number(
                    calculated, {"rounding": decimals, "text": str(calculated)}, rules
                )
            changes.append({
                "shape_id": shape_id,
                "slide": int(claim.get("slide", 1)),
                "source_value": source_display,
                "override_value": display,
                "reason": str(override["reason"]),
                "author": str(override["author"]),
                "at": str(override["at"]),
            })
        elif expected != display:
            issues.append(Issue("calc.source_manifest", int(claim.get("slide", 1)), shape_id,
                                f"source={expected!r}, manifest={display!r}"))
    if pptx_path is not None:
        issues.extend(check_chart_series(pptx_path, claims, workbooks, resolved_sources))
    for workbook in workbooks.values():
        workbook.close()
    token_issues, token_warnings = check_numeric_tokens(prs, rules, payload)
    issues.extend(token_issues)
    warnings.extend(token_warnings)
    return issues, changes, warnings


def check_numeric_tokens(prs, rules, manifest):
    config = rules["numeric_tokens"]
    token_re = re.compile(config["pattern"])
    context_patterns = [re.compile(item["pattern"]) for item in config["global_text_whitelist"]]
    job_allowed = {}
    for item in manifest.get("token_whitelist", []):
        job_allowed[(int(item["slide"]), str(item["token"]))] = item

    registered = set()
    for claim in manifest["claims"]:
        for placement in claim["placements"]:
            registered.update((int(placement["slide"]), match.group(0))
                              for match in token_re.finditer(str(placement["text"])))

    issues, emitted, whitelist_uses = [], set(), {}
    for page, slide in enumerate(prs.slides, 1):
        containers = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                containers.append((shape.name, text(shape).strip()))
            if getattr(shape, "has_table", False):
                for row_index, row in enumerate(shape.table.rows):
                    for col_index, cell in enumerate(row.cells):
                        containers.append((f"{shape.name}[{row_index},{col_index}]", cell.text.strip()))
        for location, full_text in containers:
            for match in token_re.finditer(full_text):
                token = match.group(0)
                if (page, token) in registered:
                    continue
                if (page, token) in job_allowed:
                    whitelist_uses.setdefault((page, token), []).append(location)
                    continue
                if any(pattern.fullmatch(token) for pattern in context_patterns):
                    continue
                key = (page, location, token)
                if key not in emitted:
                    emitted.add(key)
                    issues.append(Issue("claim.unregistered_numeric_token", page, location,
                                        f"등록되지 않은 숫자 토큰: {token!r}"))
    warnings = []
    for (page, token), locations in sorted(whitelist_uses.items()):
        item = job_allowed[(page, token)]
        unique_locations = sorted(set(locations))
        warnings.append(Issue(
            "token.whitelist_used", page, ", ".join(unique_locations),
            f"token={token!r}, uses={len(locations)}, reason={item['reason']!r}",
        ))
    return issues, warnings


def check_shape_placement(shape, placement, rules, page, shape_id, issues):
    unit = rules["units"]["emu_per_inch"]
    tolerance = rules["units"]["bounds_tolerance_emu"]
    actual = {"x": shape.left, "y": shape.top, "w": shape.width, "h": shape.height}
    for key, value in placement["bounds"].items():
        expected = round(float(value) * unit)
        if abs(actual[key] - expected) > tolerance:
            issues.append(Issue("claim.source_manifest_pptx", page, shape_id,
                                f"bounds.{key} manifest={expected}EMU, pptx={actual[key]}EMU"))
    runs = [run for paragraph in shape.text_frame.paragraphs for run in paragraph.runs if run.text]
    spec = placement["font"]
    if spec.get("face") and any(run.font.name != spec["face"] for run in runs):
        issues.append(Issue("claim.source_manifest_pptx", page, shape_id, "font.face 불일치"))
    if spec.get("size") is not None and any(
            run.font.size is None or abs(run.font.size.pt - float(spec["size"])) > 0.01 for run in runs):
        issues.append(Issue("claim.source_manifest_pptx", page, shape_id, "font.size 불일치"))
    if any(bool(run.font.bold) != bool(spec.get("bold")) for run in runs):
        issues.append(Issue("claim.source_manifest_pptx", page, shape_id, "font.bold 불일치"))
    actual_align = ALIGN_NAMES.get(shape.text_frame.paragraphs[0].alignment)
    if actual_align != placement.get("align", "left"):
        issues.append(Issue("claim.source_manifest_pptx", page, shape_id,
                            f"align manifest={placement.get('align')}, pptx={actual_align}"))
    anchor = str(shape.text_frame.vertical_anchor).split()[0].lower()
    anchor_map = {"middle": "middle", "top": "top", "bottom": "bottom"}
    actual_valign = anchor_map.get(anchor, anchor)
    if actual_valign != placement.get("valign", "middle"):
        issues.append(Issue("claim.source_manifest_pptx", page, shape_id,
                            f"valign manifest={placement.get('valign')}, pptx={actual_valign}"))


def audit(path, rules, manifest_path=None, source_root=None):
    rules = style_rules(rules, manifest_path)
    prs = Presentation(str(path))
    checks = (check_fonts, check_notation, check_negative_red, check_red_runs_per_line, check_font_sizes,
              check_table_alignments, check_footnotes,
              check_overflow, check_title_right, check_table_geometry,
              check_canvas_and_content, check_chip_geometry)
    issues = check_preflight_alignment(rules)
    issues.extend(issue for check in checks for issue in check(prs, rules))
    changes, warnings = [], []
    if manifest_path:
        claim_issues, changes, warnings = check_claims(
            prs, rules, manifest_path, source_root, path
        )
        issues.extend(claim_issues)
    issues.sort(key=lambda item: (item.slide, item.rule, item.shape, item.evidence))
    warnings.sort(key=lambda item: (item.slide, item.rule, item.shape, item.evidence))
    return {"file": path.name, "status": "FAIL" if issues else "PASS",
            "issues": [asdict(issue) for issue in issues],
            "warnings": [asdict(warning) for warning in warnings], "changes": changes}


def audit_safe(path, rules, manifest_path=None, source_root=None):
    try:
        return audit(path, rules, manifest_path, source_root)
    except Exception as error:
        return {
            "file": path.name,
            "status": "ERROR",
            "issues": [],
            "warnings": [],
            "changes": [],
            "error": f"{type(error).__name__}: {error}",
        }


def verify(results, expected_path):
    payload = json.loads(expected_path.read_text(encoding="utf-8"))
    expected = {payload["golden"]["file"]: payload["golden"].get("static_expected", "PASS")}
    expected.update({item["file"]: item.get("static_expected", item["expected"])
                     for item in payload["fixtures"]})
    actual = {item["file"]: item["status"] for item in results}
    return [f"{name}: expected {status}, got {actual.get(name, 'MISSING')}"
            for name, status in expected.items()
            if status != "DEFERRED" and actual.get(name) != status]


def path_default_manifest(pptx_path):
    candidate = pptx_path.with_name("manifest.json")
    return candidate if candidate.exists() else None


def aggregate_status(results):
    priority = {"PASS": 0, "SKIP": 1, "FAIL": 2, "ERROR": 3}
    return max((item["status"] for item in results), key=priority.__getitem__)


def main(argv=None):
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("target", type=Path)
    parser.add_argument("--rules", type=Path, default=Path(__file__).with_name("house-rules.yaml"))
    parser.add_argument("--json", action="store_true")
    parser.add_argument("--manifest", type=Path)
    parser.add_argument("--source-root", type=Path)
    args = parser.parse_args(argv)
    paths = sorted(args.target.glob("*.pptx")) if args.target.is_dir() else [args.target]
    if not paths:
        error = f"no pptx files found: {args.target}"
        if args.json:
            print(json.dumps({
                "status": "ERROR",
                "results": [{
                    "file": str(args.target),
                    "status": "ERROR",
                    "issues": [],
                    "warnings": [],
                    "changes": [],
                    "error": error,
                }],
                "expected_mismatches": [],
            }, ensure_ascii=False, indent=2))
        else:
            print(error, file=sys.stderr)
        return 2
    rules = load_rules(args.rules)
    expected_path = args.target / "expected_results.json" if args.target.is_dir() else None
    verifying_fixtures = bool(expected_path and expected_path.exists())
    manifest_by_file = {}
    if verifying_fixtures:
        expected_payload = json.loads(expected_path.read_text(encoding="utf-8"))
        manifest_by_file = {
            item["file"]: args.target / item["manifest"]
            for item in expected_payload["fixtures"] if item.get("manifest")
        }
        golden = expected_payload["golden"]
        if golden.get("manifest"):
            manifest_by_file[golden["file"]] = args.target / golden["manifest"]
    elif args.manifest:
        manifest_by_file[paths[0].name] = args.manifest
    elif len(paths) == 1:
        default_manifest = path_default_manifest(paths[0])
        if default_manifest:
            manifest_by_file[paths[0].name] = default_manifest
    results = [audit_safe(path, rules, manifest_by_file.get(path.name), args.source_root)
               for path in paths]
    mismatches = verify(results, expected_path) if verifying_fixtures else []
    payload_status = ("PASS" if not mismatches else "FAIL") if verifying_fixtures else aggregate_status(results)
    payload = {
        "status": payload_status,
        "results": results,
        "expected_mismatches": mismatches,
    }
    if verifying_fixtures:
        payload["fixture_match"] = not mismatches
    if args.json:
        print(json.dumps(payload, ensure_ascii=False, indent=2))
    else:
        for result in results:
            print(f"{result['status']:4} {result['file']}")
            if result.get("error"):
                print(f"  audit.error: {result['error']}")
            for issue in result["issues"]:
                print(f"  p{issue['slide']} {issue['rule']}: {issue['evidence']}")
            for warning in result.get("warnings", []):
                print(f"  WARNING p{warning['slide']} {warning['rule']}: {warning['evidence']}")
            for change in result.get("changes", []):
                print(f"  CHANGE p{change['slide']} {change['shape_id']}: "
                      f"{change['source_value']!r} -> {change['override_value']!r} "
                      f"({change['reason']}, {change['author']}, {change['at']})")
        if verifying_fixtures:
            print("EXPECTED MATCH" if not mismatches else "EXPECTED MISMATCH")
            for mismatch in mismatches:
                print(f"  {mismatch}")
    if verifying_fixtures:
        return 0 if not mismatches else 1
    if any(result["status"] == "ERROR" for result in results):
        return 2
    return 1 if any(result["status"] == "FAIL" for result in results) else 0


if __name__ == "__main__":
    raise SystemExit(main())
