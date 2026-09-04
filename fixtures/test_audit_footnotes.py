import unittest

import yaml
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from audit import check_footnotes


class FootnoteTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        with open("house-rules.yaml", encoding="utf-8") as handle:
            all_rules = yaml.safe_load(handle)
        cls.rules = dict(all_rules)
        cls.rules.update(all_rules["styles"]["shin-ppt1"])

    def presentation(self, notes=None, *, y=None, name="footer/text", rule=True,
                     rule_thickness=None):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        if notes is None:
            return presentation
        zones = self.rules["zones"]
        if y is None:
            y = zones["footnote_bottom_y"] - zones["footnote_line_step"] * len(notes)
        shape = slide.shapes.add_textbox(Inches(0.65), Inches(y), Inches(10), Inches(0.9))
        shape.name = name
        shape.text_frame.clear()
        for index, note in enumerate(notes):
            paragraph = shape.text_frame.paragraphs[0] if index == 0 else shape.text_frame.add_paragraph()
            paragraph.text = note
        if rule:
            thickness = (self.rules["components"]["footnote_rule_thickness"]
                         if rule_thickness is None else rule_thickness)
            divider = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(y - 0.12), Inches(10), Inches(thickness)
            )
            divider.name = "footer/rule"
        return presentation

    def test_missing_footnote_is_skipped(self):
        self.assertEqual(check_footnotes(self.presentation(), self.rules), [])

    def test_four_lines_at_rule_position_pass(self):
        notes = ["※ 원천", "* 주석 2", "* 주석 3", "* 주석 4"]
        self.assertEqual(check_footnotes(self.presentation(notes), self.rules), [])

    def test_five_lines_fail_maximum(self):
        notes = ["※ 원천", "* 주석 2", "* 주석 3", "* 주석 4", "* 주석 5"]
        issues = check_footnotes(self.presentation(notes), self.rules)
        self.assertEqual([issue.rule for issue in issues], ["zones.footnote_max_lines"])

    def test_wrong_vertical_position_fails(self):
        issues = check_footnotes(self.presentation(["※ 원천"], y=7.5), self.rules)
        self.assertEqual([issue.rule for issue in issues], ["zones.footnote_bottom_y"])

    def test_missing_rule_with_footnote_fails(self):
        issues = check_footnotes(self.presentation(["※ 원천"], rule=False), self.rules)
        self.assertEqual([issue.rule for issue in issues], ["components.footnote_rule_thickness"])

    def test_wrong_rule_thickness_fails(self):
        issues = check_footnotes(
            self.presentation(["※ 원천"], rule_thickness=0.008), self.rules
        )
        self.assertEqual([issue.rule for issue in issues], ["components.footnote_rule_thickness"])

    def test_page_number_footer_is_not_a_note(self):
        self.assertEqual(check_footnotes(self.presentation(["8"]), self.rules), [])


if __name__ == "__main__":
    unittest.main()
