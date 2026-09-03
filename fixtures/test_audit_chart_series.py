import hashlib
import tempfile
import unittest
from pathlib import Path

from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from audit import check_chart_series


class ChartSeriesTests(unittest.TestCase):
    def make_inputs(self, root: Path, chart_values):
        source_path = root / "source.xlsx"
        source = Workbook()
        sheet = source.active
        sheet.title = "시계열"
        sheet.append([17.2, 0])
        sheet.append([18.9, 22.9])
        source.save(source_path)

        deck_path = root / "chart.pptx"
        deck = Presentation()
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        data = ChartData()
        data.categories = [1, 2]  # 축 값은 원천에 없어도 검사 대상이 아니다.
        data.add_series("거래대금", chart_values)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, Inches(1), Inches(1), Inches(5), Inches(3), data
        )
        chart.name = "chartLine/chart"
        deck.save(deck_path)
        return source_path, deck_path

    def inspect(self, chart_values, register_range=False):
        with tempfile.TemporaryDirectory() as temp:
            source_path, deck_path = self.make_inputs(Path(temp), chart_values)
            workbook = load_workbook(source_path, data_only=True, read_only=True)
            try:
                claims = [{
                    "source": {
                        "sheet": "시계열",
                    },
                    "placements": [{"slide": 1}],
                }]
                chart_series = []
                if register_range:
                    chart_series.append({
                        "slide": 1,
                        "chart": "chartLine/chart",
                        "series": 1,
                        "name": "거래대금",
                        "source": {
                            "file": "source.xlsx",
                            "file_hash": "sha256:" + hashlib.sha256(source_path.read_bytes()).hexdigest(),
                            "sheet": "시계열",
                            "ref": "A1:A2",
                        },
                    })
                return check_chart_series(
                    deck_path, claims, {source_path: workbook}, {id(claims[0]): source_path},
                    chart_series, Path(temp),
                )
            finally:
                workbook.close()

    def test_source_backed_series_values_pass(self):
        issues, warnings = self.inspect([17.2, 18.9])
        self.assertEqual(issues, [])
        self.assertEqual([warning.rule for warning in warnings], [
            "claim.chart_series_range_unverified",
        ])

    def test_unregistered_series_value_fails_but_axis_values_do_not(self):
        issues, _ = self.inspect([99.9, 18.9])
        self.assertEqual(len(issues), 1)
        self.assertEqual(issues[0].rule, "claim.unregistered_chart_series_value")
        self.assertIn("99.9", issues[0].evidence)

    def test_registered_range_compares_values_in_order(self):
        issues, warnings = self.inspect([17.2, 18.9], register_range=True)
        self.assertEqual(issues, [])
        self.assertEqual(warnings, [])

    def test_registered_range_rejects_value_from_other_column(self):
        issues, warnings = self.inspect([17.2, 22.9], register_range=True)
        self.assertEqual([issue.rule for issue in issues], [
            "claim.unregistered_chart_series_value",
        ])
        self.assertEqual(warnings, [])
        self.assertIn("2: chart=22.9, source=18.9", issues[0].evidence)

    def test_registered_range_rejects_reordered_values(self):
        issues, _ = self.inspect([18.9, 17.2], register_range=True)
        self.assertEqual([issue.rule for issue in issues], [
            "claim.unregistered_chart_series_value",
        ])


if __name__ == "__main__":
    unittest.main()
