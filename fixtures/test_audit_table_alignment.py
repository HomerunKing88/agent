import unittest

import yaml
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

from audit import check_table_alignments


class TableAlignmentTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        with open("house-rules.yaml", encoding="utf-8") as handle:
            all_rules = yaml.safe_load(handle)
        cls.rules = dict(all_rules)
        cls.rules.update(all_rules["styles"]["corporate-strategy-ppt"])

    def make_table(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        table = slide.shapes.add_table(3, 3, 0, 0, 3000000, 1000000).table
        values = [
            ["제품", "구분", "금액"],
            ["퉁퉁이별티", "온라인", "2,335,436"],
            ["유자", "10", "621,565"],
        ]
        for row_index, row in enumerate(values):
            for col_index, value in enumerate(row):
                cell = table.cell(row_index, col_index)
                cell.text = value
                paragraph = cell.text_frame.paragraphs[0]
                if row_index == 0:
                    paragraph.alignment = PP_ALIGN.CENTER
                elif col_index == 0:
                    paragraph.alignment = PP_ALIGN.LEFT
                elif value.replace(",", "").isdigit():
                    paragraph.alignment = PP_ALIGN.RIGHT
                else:
                    paragraph.alignment = PP_ALIGN.CENTER
        return presentation, table

    def test_header_numeric_default_and_text_column_alignments_pass(self):
        presentation, _ = self.make_table()
        self.assertEqual(check_table_alignments(presentation, self.rules), [])

    def test_centered_numeric_cell_fails(self):
        presentation, table = self.make_table()
        table.cell(1, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        issues = check_table_alignments(presentation, self.rules)
        self.assertEqual([issue.rule for issue in issues], ["table.numeric_col_align"])

    def test_mixed_column_text_uses_default_alignment(self):
        presentation, table = self.make_table()
        table.cell(1, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        issues = check_table_alignments(presentation, self.rules)
        self.assertEqual([issue.rule for issue in issues], ["table.default_align"])

    def test_bold_unfilled_subtotal_may_use_default_alignment(self):
        presentation, table = self.make_table()
        for cell in table.rows[2].cells:
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.bold = True
        table.cell(2, 2).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        self.assertEqual(check_table_alignments(presentation, self.rules), [])


if __name__ == "__main__":
    unittest.main()
