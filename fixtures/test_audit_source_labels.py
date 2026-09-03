import unittest

from openpyxl import Workbook
from pptx import Presentation

from audit import check_source_label, missing_label_ref_warnings


class SourceLabelTests(unittest.TestCase):
    def presentation(self, labels):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        shape = slide.shapes.add_table(len(labels), 2, 0, 0, 3000000, 1000000)
        shape.name = "table/perf"
        for row, label in enumerate(labels):
            shape.table.cell(row, 0).text = label
            shape.table.cell(row, 1).text = str(row)
        return presentation

    def workbook(self):
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "실적"
        sheet["B17"] = "퉁퉁이별티"
        sheet["B18"] = "퉁퉁이별티 유자"
        sheet["B19"] = "퉁퉁이별티 호박팥"
        sheet["B20"] = "퉁퉁이별티 카라멜 루이보스"
        return workbook

    def claim(self, *, placement_type="cell", label_ref="B18", row=1, shape_id="VALUE"):
        placement = {"slide": 1, "type": placement_type}
        if placement_type == "cell":
            placement.update({"table": "table/perf", "row": row, "col": 1, "text": "43"})
        else:
            placement.update({"name": "CLAIM", "text": "43"})
        source = {"file": "source.xlsx", "sheet": "실적", "ref": "C18"}
        if label_ref is not None:
            source["label_ref"] = label_ref
        return {"shape_id": shape_id, "source": source, "placements": [placement]}

    def catalog(self):
        return {("source.xlsx", "실적", 1, "table/perf"): {
            "퉁퉁이별티", "퉁퉁이별티 유자", "퉁퉁이별티 호박팥",
            "퉁퉁이별티 카라멜 루이보스",
        }}

    def test_shortened_deck_label_passes(self):
        issues = check_source_label(
            self.presentation(["제품", "유자"]), self.claim(), self.workbook(), set(), self.catalog()
        )
        self.assertEqual(issues, [])

    def test_different_row_label_fails(self):
        issues = check_source_label(
            self.presentation(["제품", "호박팥"]), self.claim(), self.workbook(), set(), self.catalog()
        )
        self.assertEqual([issue.rule for issue in issues], ["claim.source_label_mismatch"])

    def test_shape_placement_and_missing_label_ref_skip(self):
        presentation = self.presentation(["제품", "다른 값"])
        workbook = self.workbook()
        self.assertEqual(check_source_label(
            presentation, self.claim(placement_type="shape"), workbook, set(), self.catalog()
        ), [])
        self.assertEqual(check_source_label(
            presentation, self.claim(label_ref=None), workbook, set(), self.catalog()
        ), [])

    def test_same_source_row_is_checked_once(self):
        presentation = self.presentation(["제품", "호박팥"])
        workbook = self.workbook()
        checked = set()
        first = check_source_label(presentation, self.claim(), workbook, checked, self.catalog())
        second = check_source_label(presentation, self.claim(), workbook, checked, self.catalog())
        self.assertEqual(len(first), 1)
        self.assertEqual(second, [])

    def test_common_prefix_is_not_an_abbreviation(self):
        issues = check_source_label(
            self.presentation(["제품", "퉁퉁이별티"]), self.claim(),
            self.workbook(), set(), self.catalog()
        )
        self.assertEqual([issue.rule for issue in issues], ["claim.source_label_mismatch"])

    def test_swapped_first_two_rows_both_fail(self):
        presentation = self.presentation(["제품", "유자", "퉁퉁이별티"])
        workbook = self.workbook()
        checked = set()
        issues = check_source_label(
            presentation, self.claim(label_ref="B17", row=1), workbook, checked, self.catalog()
        )
        issues.extend(check_source_label(
            presentation, self.claim(label_ref="B18", row=2), workbook, checked, self.catalog()
        ))
        self.assertEqual([issue.rule for issue in issues], [
            "claim.source_label_mismatch", "claim.source_label_mismatch",
        ])

    def test_missing_label_refs_are_grouped_into_one_warning(self):
        claims = [
            self.claim(label_ref=None, shape_id="VALUE_A"),
            self.claim(label_ref=None, shape_id="VALUE_B"),
            self.claim(label_ref=None, placement_type="shape", shape_id="SHAPE_VALUE"),
        ]
        warnings = missing_label_ref_warnings(claims)
        self.assertEqual([warning.rule for warning in warnings], [
            "claim.source_label_unverified",
        ])
        self.assertIn("claim=2", warnings[0].evidence)


if __name__ == "__main__":
    unittest.main()
