import unittest

from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from render_check import (
    SvgTextRender,
    inspect_libreoffice_presentation,
    inspect_presentation,
    normalized_font,
)


class Collection:
    def __init__(self, items):
        self.items = items
        self.Count = len(items)

    def Item(self, index):
        return self.items[index - 1]


class CallableCollection(Collection):
    def __call__(self, *args):
        return self


class Font:
    Name = "맑은 고딕"


class TextRange:
    Text = "한 줄이 자동으로 줄바꿈되는 문장"
    BoundLeft = 0
    BoundTop = 590
    BoundWidth = 40
    BoundHeight = 20
    Font = Font()

    def Runs(self, *args):
        if args:
            return self
        return CallableCollection([self])

    def Lines(self, *args):
        return CallableCollection([object(), object()])


class Frame:
    HasText = -1
    TextRange = TextRange()


class Shape:
    Type = 1
    HasTextFrame = -1
    HasTable = 0
    TextFrame2 = Frame()
    Name = "fixture/overflow"
    Left = 0
    Top = 590
    Width = 40
    Height = 10


class Slide:
    Shapes = Collection([Shape()])


class Presentation:
    Slides = Collection([Slide()])


class RenderCheckTest(unittest.TestCase):
    def test_overflow_and_page_bottom_are_reported_but_wrap_is_allowed(self):
        rules = {
            "fonts": {"heading": "HY헤드라인M"},
            "units": {"pt_per_inch": 72},
            "qa": {"canvas_overflow_tolerance_in": 0.01, "text_max_ymax_pt": 593},
        }
        issues, skips = inspect_presentation(Presentation(), rules, missing_heading=False)
        self.assertEqual(skips, [])
        self.assertEqual(
            {issue.rule for issue in issues},
            {"render.text_overflow", "render.page_text_ymax"},
        )

    def test_libreoffice_checks_available_font_and_skips_only_untrusted_fonts(self):
        presentation = PptxPresentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])

        body = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        body.name = "body"
        body.text_frame.paragraphs[0].add_run().text = "본문"
        body.text_frame.paragraphs[0].runs[0].font.name = "맑은 고딕"

        heading = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(1))
        heading.name = "heading"
        heading.text_frame.paragraphs[0].add_run().text = "제목"
        heading.text_frame.paragraphs[0].runs[0].font.name = "HY헤드라인M"

        substituted = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(2), Inches(1))
        substituted.name = "substituted"
        substituted.text_frame.paragraphs[0].add_run().text = "대체"
        substituted.text_frame.paragraphs[0].runs[0].font.name = "맑은 고딕"

        rendered = [[
            # Three rendered lines without an expanded bound are normal wrapping.
            SvgTextRender("본문", 72, 72, 216, 144, 3, frozenset({"맑은 고딕"})),
            SvgTextRender("제목", 72, 144, 216, 216, 1, frozenset({"HY헤드라인M"})),
            SvgTextRender("대체", 72, 216, 216, 288, 1, frozenset({"Arial"})),
        ]]
        rules = {
            "units": {"pt_per_inch": 72},
            "qa": {"canvas_overflow_tolerance_in": 0.01, "text_max_ymax_pt": 593},
        }
        issues, skips, inspected = inspect_libreoffice_presentation(
            presentation, rendered, rules, {normalized_font("맑은 고딕")},
        )

        self.assertEqual(issues, [])
        self.assertEqual(inspected, 1)
        self.assertEqual({skip["shape"] for skip in skips}, {"heading", "substituted"})
        self.assertTrue(any("font missing" in skip["reason"] for skip in skips))
        self.assertTrue(any("font substitution detected" in skip["reason"] for skip in skips))


if __name__ == "__main__":
    unittest.main()
