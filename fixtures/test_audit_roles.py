import unittest

import yaml
from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement

from audit import check_font_sizes, minimum_font_size


class MinimumFontSizeTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        with open("house-rules.yaml", encoding="utf-8") as handle:
            cls.rules = yaml.safe_load(handle)["styles"]["corporate-strategy-ppt"]

    def test_named_roles_use_house_rule_mapping(self):
        cases = {
            "icon_badge/glyph": "icon_badge_glyph_pt",
            "stack100/seg_label#2": "chart_value_label_pt",
            "waterfall/tick_label": "chart_axis_label_pt",
            "waterfall/value": "waterfall_value_pt",
        }
        for shape_name, size_key in cases.items():
            with self.subTest(shape_name=shape_name):
                self.assertEqual(
                    minimum_font_size(shape_name, self.rules),
                    float(self.rules["sizes"][size_key]),
                )

    def test_claim_fallback_and_unknown_structure_error(self):
        self.assertEqual(
            minimum_font_size("CLAIM_REVENUE", self.rules),
            float(self.rules["sizes"][self.rules["role_min_pt"]["_claim_shape"]]),
        )
        with self.assertRaises(ValueError):
            minimum_font_size("future_helper/text", self.rules)

    def test_table_cells_keep_table_minimum(self):
        self.assertEqual(
            minimum_font_size("table/perf", self.rules, is_table=True),
            float(self.rules["sizes"]["table_body_min_pt"]),
        )

    def table_with_unset_run_size(self, paragraph_size=None):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        table_shape = slide.shapes.add_table(1, 1, 0, 0, 2000000, 500000)
        table_shape.name = "table/perf"
        paragraph = table_shape.table.cell(0, 0).text_frame.paragraphs[0]
        paragraph.text = "값"
        if paragraph_size is not None:
            properties = paragraph._p.get_or_add_pPr()
            default_run = OxmlElement("a:defRPr")
            default_run.set("sz", str(paragraph_size * 100))
            properties.append(default_run)
        return presentation

    def test_unset_run_size_is_resolved_and_reported(self):
        warnings = []
        issues = check_font_sizes(self.table_with_unset_run_size(), self.rules, warnings)
        self.assertEqual(issues, [])
        self.assertEqual([warning.rule for warning in warnings], ["sizes.font_size_inherited"])
        self.assertIn("master otherStyle", warnings[0].evidence)

    def test_inherited_small_size_fails_minimum(self):
        warnings = []
        issues = check_font_sizes(self.table_with_unset_run_size(6), self.rules, warnings)
        self.assertEqual([issue.rule for issue in issues], ["sizes.body_min_pt"])
        self.assertEqual([warning.rule for warning in warnings], ["sizes.font_size_inherited"])


if __name__ == "__main__":
    unittest.main()
