import unittest

import yaml
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches

from audit import DRAWINGML_NS, check_component_text_and_chip_rules


class ComponentRuleTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        with open("house-rules.yaml", encoding="utf-8") as handle:
            all_rules = yaml.safe_load(handle)
        cls.rules = dict(all_rules)
        cls.rules.update(all_rules["styles"]["corporate-strategy-ppt"])
        cls.shin_rules = dict(all_rules)
        cls.shin_rules.update(all_rules["styles"]["shin-ppt1"])

    def presentation(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        marker = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(0.2), Inches(0.26))
        marker.name = "bullets/marker"
        marker.text = self.rules["components"]["bullet_marker"]
        line = slide.shapes.add_textbox(Inches(1.2), Inches(1), Inches(4), Inches(0.4))
        line.name = "bullets/line"
        line.text = "일반 불릿"
        conclusion = slide.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(4), Inches(0.4))
        conclusion.name = "bullets/line#2"
        conclusion.text = self.rules["components"]["conclusion_prefix"] + " 결론"
        caption = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(0.3))
        caption.name = "table/caption"
        caption.text = self.rules["components"]["table_caption_prefix"] + "표 설명"
        chip = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2.5), Inches(1.95), Inches(0.3)
        )
        chip.name = "chip/bg"
        return presentation

    def test_valid_components_pass(self):
        self.assertEqual(check_component_text_and_chip_rules(self.presentation(), self.rules), [])

    def test_wrong_marker_caption_and_conclusion_prefix_fail(self):
        presentation = self.presentation()
        slide = presentation.slides[0]
        next(shape for shape in slide.shapes if shape.name == "bullets/marker").text = "•"
        next(shape for shape in slide.shapes if shape.name == "bullets/line#2").text = "결론"
        next(shape for shape in slide.shapes if shape.name == "table/caption").text = "표 설명"
        issues = check_component_text_and_chip_rules(presentation, self.rules)
        self.assertEqual({issue.rule for issue in issues}, {
            "components.bullet_marker",
            "components.conclusion_prefix",
            "components.table_caption_prefix",
        })

    def test_wrong_chip_radius_fails(self):
        presentation = self.presentation()
        chip = next(shape for shape in presentation.slides[0].shapes if shape.name == "chip/bg")
        adjustment = chip._element.find(
            f".//{{{DRAWINGML_NS}}}gd[@name='adj']"
        )
        if adjustment is None:
            adjustments = chip._element.find(f".//{{{DRAWINGML_NS}}}avLst")
            adjustment = OxmlElement("a:gd")
            adjustment.set("name", "adj")
            adjustments.append(adjustment)
        adjustment.set("fmla", "val 33333")
        issues = check_component_text_and_chip_rules(presentation, self.rules)
        self.assertEqual([issue.rule for issue in issues], ["components.chip.radius"])

    def test_missing_style_rules_skip(self):
        presentation = self.presentation()
        self.assertEqual(
            check_component_text_and_chip_rules(presentation, self.shin_rules), []
        )


if __name__ == "__main__":
    unittest.main()
