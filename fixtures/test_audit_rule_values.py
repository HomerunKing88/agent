import copy
import unittest

import yaml
from pptx import Presentation
from pptx.util import Inches, Pt

from audit import check_canvas_and_content, check_fonts


class AuditRuleValueTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        with open("house-rules.yaml", encoding="utf-8") as handle:
            all_rules = yaml.safe_load(handle)
        cls.corporate = dict(all_rules)
        cls.corporate.update(all_rules["styles"]["corporate-strategy-ppt"])
        cls.shin = dict(all_rules)
        cls.shin.update(all_rules["styles"]["shin-ppt1"])

    @staticmethod
    def three_font_presentation():
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        for index, font in enumerate(("HY헤드라인M", "맑은 고딕", "Courier New")):
            shape = slide.shapes.add_textbox(Inches(1), Inches(1 + index), Inches(2), Inches(0.4))
            shape.name = f"font/{index}"
            run = shape.text_frame.paragraphs[0].add_run()
            run.text = font
            run.font.name = font
            run.font.size = Pt(10)
        return presentation

    @staticmethod
    def draft_tag_presentation():
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(8.4), Inches(8), Inches(2.65), Inches(0.2))
        shape.name = "header/draft_tag"
        shape.text = "DRAFT"
        return presentation

    def test_allowed_count_controls_actual_font_count(self):
        presentation = self.three_font_presentation()
        self.assertEqual([issue.rule for issue in check_fonts(presentation, self.corporate)],
                         ["forbidden.third_font"])
        rules = copy.deepcopy(self.corporate)
        rules["fonts"]["allowed_count"] = 3
        self.assertEqual(check_fonts(presentation, rules), [])

    def test_content_exemption_uses_rule_role_name(self):
        presentation = self.draft_tag_presentation()
        self.assertEqual(check_canvas_and_content(presentation, self.corporate), [])
        rules = copy.deepcopy(self.corporate)
        rules["zones"]["content_max_y_exempt"] = []
        self.assertEqual(
            [issue.rule for issue in check_canvas_and_content(presentation, rules)],
            ["zones.content_max_y"],
        )

    def test_missing_exemption_key_means_no_exemption(self):
        presentation = self.draft_tag_presentation()
        self.assertEqual(
            [issue.rule for issue in check_canvas_and_content(presentation, self.shin)],
            ["zones.content_max_y"],
        )


if __name__ == "__main__":
    unittest.main()
