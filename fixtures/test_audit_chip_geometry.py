import unittest

import yaml
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from audit import check_chip_geometry


class ChipGeometryTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        with open("house-rules.yaml", encoding="utf-8") as handle:
            all_rules = yaml.safe_load(handle)
        cls.rules = dict(all_rules)
        cls.rules.update(all_rules["styles"]["corporate-strategy-ppt"])
        cls.shin_rules = dict(all_rules)
        cls.shin_rules.update(all_rules["styles"]["shin-ppt1"])

    def presentation(self, chips):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        for index, (x, y, bg_width, desc_width) in enumerate(chips, 1):
            bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(bg_width), Inches(0.30),
            )
            bg.name = "chip/bg" if index == 1 else f"chip/bg#{index}"
            desc = slide.shapes.add_textbox(
                Inches(x + bg_width + 0.10), Inches(y + 0.04),
                Inches(desc_width), Inches(0.22),
            )
            desc.name = "chip/desc" if index == 1 else f"chip/desc#{index}"
        return presentation

    def test_full_and_margin_clipped_widths_pass(self):
        presentation = self.presentation([
            (0.65, 2.10, 1.95, 5.20),
            (6.00, 2.10, 1.95, 2.99),
        ])
        self.assertEqual(check_chip_geometry(presentation, self.rules), [])

    def test_desc_beyond_right_margin_fails(self):
        presentation = self.presentation([(6.00, 2.10, 1.95, 5.20)])
        issues = check_chip_geometry(presentation, self.rules)
        self.assertEqual([issue.rule for issue in issues], ["layout.chip_desc_canvas"])

    def test_width_above_cap_fails(self):
        presentation = self.presentation([(0.65, 2.10, 1.95, 5.30)])
        issues = check_chip_geometry(presentation, self.rules)
        self.assertEqual([issue.rule for issue in issues], ["layout.chip_desc_width"])

    def test_missing_chips_skip_style_without_chip_rules(self):
        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[6])
        self.assertEqual(check_chip_geometry(presentation, self.shin_rules), [])


if __name__ == "__main__":
    unittest.main()
