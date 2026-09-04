#!/usr/bin/env python3
"""Render checks for real text bounds and wrapping.

Windows uses PowerPoint COM. Other platforms use LibreOffice's headless SVG
export, whose ``BoundingBox`` expands when rendered text escapes its shape.
"""
from __future__ import annotations

import argparse
import json
import platform
import re
import shutil
import struct
import subprocess
import sys
import tempfile
import xml.etree.ElementTree as ET
from dataclasses import asdict, dataclass
from functools import lru_cache
from pathlib import Path

import yaml
from pptx import Presentation


@dataclass(frozen=True)
class RenderIssue:
    rule: str
    slide: int
    shape: str
    evidence: str


@dataclass(frozen=True)
class SvgTextRender:
    text: str
    left: float
    top: float
    right: float
    bottom: float
    lines: int
    fonts: frozenset[str]


@dataclass(frozen=True)
class PptxTextTarget:
    name: str
    text: str
    left: int
    top: int
    width: int
    height: int
    text_frame: object


def load_rules(path: Path) -> dict:
    with path.open(encoding="utf-8") as handle:
        rules = yaml.safe_load(handle)
    for section in ("units", "styles", "default_style"):
        if section not in rules:
            raise ValueError(f"house-rules missing section: {section}")
    return rules


def style_rules(rules: dict, style: str | None = None, manifest_path: Path | None = None) -> dict:
    """Select render rules without silently falling back to a default style."""
    if manifest_path is not None:
        payload = json.loads(manifest_path.read_text(encoding="utf-8"))
        manifest_style = payload.get("style")
        if not manifest_style:
            raise ValueError("manifest style is required")
        if style is not None and style != manifest_style:
            raise ValueError(f"style {style!r} disagrees with manifest style {manifest_style!r}")
        style = manifest_style
    if not style:
        raise ValueError("style is unknown: provide --manifest or --style")
    selected = rules["styles"].get(style)
    if selected is None:
        raise ValueError(f"unknown style {style!r}; expected one of {sorted(rules['styles'])}")
    missing = sorted({"fonts", "qa"} - selected.keys())
    if missing:
        raise ValueError(f"style {style!r} missing sections: {', '.join(missing)}")
    effective = dict(rules)
    effective.update(selected)
    return effective


def result(file: Path, status: str, issues=(), skips=(), error=None) -> dict:
    payload = {
        "file": file.name,
        "status": status,
        "issues": [asdict(issue) for issue in issues],
        "skips": list(skips),
    }
    if error:
        payload["error"] = error
    return payload


def installed_font_names(winreg) -> set[str]:
    locations = (
        (winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\Microsoft\Windows NT\CurrentVersion\Fonts"),
        (winreg.HKEY_CURRENT_USER, r"SOFTWARE\Microsoft\Windows NT\CurrentVersion\Fonts"),
    )
    names = set()
    for hive, key_name in locations:
        try:
            with winreg.OpenKey(hive, key_name) as key:
                index = 0
                while True:
                    try:
                        name, value, _ = winreg.EnumValue(key, index)
                    except OSError:
                        break
                    names.add(str(name).casefold())
                    names.add(Path(str(value)).stem.casefold())
                    index += 1
        except OSError:
            continue
    return names


def font_available(font: str, installed: set[str]) -> bool:
    target = font.casefold()
    return any(target in candidate for candidate in installed)


def iter_shapes(collection):
    """Yield top-level, grouped, and table-cell shapes."""
    for index in range(1, collection.Count + 1):
        shape = collection.Item(index)
        yield shape
        if int(shape.Type) == 6:  # msoGroup
            yield from iter_shapes(shape.GroupItems)
        try:
            has_table = bool(shape.HasTable)
        except Exception:
            has_table = False
        if has_table:
            table = shape.Table
            for row in range(1, table.Rows.Count + 1):
                for col in range(1, table.Columns.Count + 1):
                    yield table.Cell(row, col).Shape


def shape_name(shape) -> str:
    try:
        return str(shape.Name)
    except Exception:
        return "<unnamed>"


def text_font_names(text_range) -> set[str]:
    names = set()
    try:
        for index in range(1, text_range.Runs().Count + 1):
            name = str(text_range.Runs(index).Font.Name).strip()
            if name:
                names.add(name)
    except Exception:
        name = str(text_range.Font.Name).strip()
        if name:
            names.add(name)
    return names


def normalized_text(value: str) -> str:
    return re.sub(r"\s+", " ", value).strip()


def normalized_font(value: str) -> str:
    return re.sub(r"[\s_-]+", "", value).casefold()


def _sfnt_names(data: bytes, base: int = 0) -> set[str]:
    """Read family/full/PostScript names from one TTF/OTF face."""
    names: set[str] = set()
    try:
        num_tables = struct.unpack_from(">H", data, base + 4)[0]
        name_offset = None
        for index in range(num_tables):
            record = base + 12 + index * 16
            tag, _, offset, _ = struct.unpack_from(">4sIII", data, record)
            if tag == b"name":
                name_offset = offset
                break
        if name_offset is None:
            return names
        _, count, strings = struct.unpack_from(">HHH", data, name_offset)
        for index in range(count):
            record = name_offset + 6 + index * 12
            platform_id, _, _, name_id, length, offset = struct.unpack_from(">HHHHHH", data, record)
            if name_id not in {1, 4, 6, 16, 17}:
                continue
            raw = data[name_offset + strings + offset:name_offset + strings + offset + length]
            encoding = "utf-16-be" if platform_id in {0, 3} else "mac_roman"
            value = raw.decode(encoding, errors="ignore").strip("\x00 ")
            if value:
                names.add(value)
    except (IndexError, struct.error, ValueError):
        pass
    return names


def font_file_names(path: Path) -> set[str]:
    names = {path.stem}
    try:
        data = path.read_bytes()
    except OSError:
        return names
    if data[:4] == b"ttcf":
        try:
            count = struct.unpack_from(">I", data, 8)[0]
            offsets = struct.unpack_from(f">{count}I", data, 12)
        except (struct.error, ValueError):
            return names
        for offset in offsets:
            names.update(_sfnt_names(data, offset))
    else:
        names.update(_sfnt_names(data))
    return names


@lru_cache(maxsize=1)
def installed_non_windows_fonts() -> set[str]:
    roots = [
        Path.home() / "Library/Fonts", Path("/Library/Fonts"), Path("/System/Library/Fonts"),
        Path.home() / ".local/share/fonts", Path("/usr/local/share/fonts"), Path("/usr/share/fonts"),
    ]
    names: set[str] = set()
    for root in roots:
        if not root.is_dir():
            continue
        for pattern in ("*.ttf", "*.otf", "*.ttc", "*.TTF", "*.OTF", "*.TTC"):
            for path in root.rglob(pattern):
                names.update(font_file_names(path))
    return {normalized_font(name) for name in names if name}


def non_windows_font_available(font: str, installed: set[str]) -> bool:
    return normalized_font(font) in installed


def soffice_executable() -> Path | None:
    found = shutil.which("soffice")
    if found:
        return Path(found)
    mac_path = Path("/Applications/LibreOffice.app/Contents/MacOS/soffice")
    return mac_path if mac_path.is_file() else None


def libreoffice_svg(path: Path, soffice: Path, work: Path) -> Path:
    output = work / "output"
    profile = work / "profile"
    output.mkdir()
    completed = subprocess.run([
        str(soffice), f"-env:UserInstallation={profile.resolve().as_uri()}", "--headless",
        "--convert-to", "svg:impress_svg_Export", "--outdir", str(output), str(path.resolve()),
    ], capture_output=True, text=True, timeout=120, check=False)
    svg = output / f"{path.stem}.svg"
    if completed.returncode or not svg.is_file():
        detail = (completed.stderr or completed.stdout).strip()
        raise RuntimeError(f"LibreOffice SVG export failed ({completed.returncode}): {detail}")
    return svg


def parse_svg_text(svg_path: Path, slide_width_in: float, slide_height_in: float) -> list[list[SvgTextRender]]:
    root = ET.parse(svg_path).getroot()
    view_box = [float(item) for item in root.attrib["viewBox"].split()]
    units_x = view_box[2] / slide_width_in
    units_y = view_box[3] / slide_height_in
    ns = {"svg": "http://www.w3.org/2000/svg"}
    slides: list[list[SvgTextRender]] = []
    for slide in root.findall(".//svg:g[@class='Slide']", ns):
        if slide.attrib.get("id") == "dummy-slide":
            continue
        rendered: list[SvgTextRender] = []
        for group in slide.findall(".//svg:g[@id]", ns):
            text_shape = group.find("./svg:text[@class='SVGTextShape']", ns)
            box = group.find("./svg:rect[@class='BoundingBox']", ns)
            if text_shape is None or box is None:
                continue
            paragraphs = text_shape.findall("./svg:tspan[@class='TextParagraph']", ns)
            text = " ".join(normalized_text("".join(paragraph.itertext()))
                            for paragraph in paragraphs)
            if not normalized_text(text):
                continue
            positions = text_shape.findall(".//svg:tspan[@class='TextPosition']", ns)
            fonts = frozenset(
                span.attrib["font-family"].strip()
                for span in text_shape.findall(".//svg:tspan[@font-family]", ns)
                if span.attrib.get("font-family", "").strip()
            )
            x, y = float(box.attrib["x"]), float(box.attrib["y"])
            width, height = float(box.attrib["width"]), float(box.attrib["height"])
            rendered.append(SvgTextRender(
                normalized_text(text), x / units_x * 72, y / units_y * 72,
                (x + width) / units_x * 72, (y + height) / units_y * 72,
                max(1, len(positions)), fonts,
            ))
        for table in slide.findall(".//svg:g[@class='com.sun.star.drawing.TableShape']", ns):
            for text_shape in table.findall(".//svg:text[@class='SVGTextShape']", ns):
                positions = text_shape.findall(".//svg:tspan[@class='TextPosition']", ns)
                if not positions:
                    continue
                paragraphs = text_shape.findall("./svg:tspan[@class='TextParagraph']", ns)
                text = " ".join(normalized_text("".join(paragraph.itertext()))
                                for paragraph in paragraphs)
                fonts = frozenset(
                    span.attrib["font-family"].strip()
                    for span in text_shape.findall(".//svg:tspan[@font-family]", ns)
                    if span.attrib.get("font-family", "").strip()
                )
                line_boxes = []
                for position in positions:
                    glyphs = position.findall("./svg:tspan[@font-size]", ns)
                    if not glyphs:
                        continue
                    x = float(position.attrib["x"])
                    baseline = float(position.attrib["y"])
                    size = max(float(glyph.attrib["font-size"].removesuffix("px")) for glyph in glyphs)
                    width = sum(float(glyph.attrib.get("textLength", 0)) for glyph in glyphs)
                    line_boxes.append((x, baseline - size * 0.85, x + width, baseline + size * 0.2))
                if not line_boxes or not text:
                    continue
                rendered.append(SvgTextRender(
                    text,
                    min(box[0] for box in line_boxes) / units_x * 72,
                    min(box[1] for box in line_boxes) / units_y * 72,
                    max(box[2] for box in line_boxes) / units_x * 72,
                    max(box[3] for box in line_boxes) / units_y * 72,
                    len(positions), fonts,
                ))
        slides.append(rendered)
    return slides


def iter_pptx_text_shapes(shapes):
    for shape in shapes:
        if getattr(shape, "shape_type", None) == 6:  # GROUP
            yield from iter_pptx_text_shapes(shape.shapes)
        if getattr(shape, "has_text_frame", False) and normalized_text(shape.text):
            yield shape
        if getattr(shape, "has_table", False):
            table = shape.table
            top = shape.top
            for row_index, row in enumerate(table.rows):
                left = shape.left
                for col_index, cell in enumerate(row.cells):
                    width = table.columns[col_index].width
                    if normalized_text(cell.text):
                        yield PptxTextTarget(
                            f"{shape.name}[{row_index},{col_index}]", cell.text,
                            left, top, width, row.height, cell.text_frame,
                        )
                    left += width
                top += row.height


def pptx_shape_fonts(shape) -> set[str]:
    return {
        run.font.name
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.font.name
    }


def match_render(shape, candidates: list[SvgTextRender], used: set[int]):
    matches = [(index, item) for index, item in enumerate(candidates)
               if index not in used and item.text == normalized_text(shape.text)]
    if not matches:
        return None
    left, top = shape.left / 12700, shape.top / 12700
    return min(matches, key=lambda pair: abs(pair[1].left - left) + abs(pair[1].top - top))


def inspect_libreoffice_presentation(presentation, rendered_slides, rules, installed_fonts):
    qa = rules["qa"]
    tolerance_pt = qa["canvas_overflow_tolerance_in"] * rules["units"]["pt_per_inch"]
    issues, skips, inspected = [], [], 0
    for slide_index, slide in enumerate(presentation.slides, 1):
        candidates = rendered_slides[slide_index - 1] if slide_index <= len(rendered_slides) else []
        used: set[int] = set()
        for shape in iter_pptx_text_shapes(slide.shapes):
            matched = match_render(shape, candidates, used)
            if matched is None:
                continue
            record_index, rendered = matched
            used.add(record_index)
            name = shape.name or "<unnamed>"
            requested = pptx_shape_fonts(shape)
            if not requested:
                skips.append({"slide": slide_index, "shape": name,
                              "reason": "font substitution cannot be verified: source font is inherited"})
                continue
            missing = sorted(font for font in requested
                             if not non_windows_font_available(font, installed_fonts))
            if missing:
                skips.append({"slide": slide_index, "shape": name,
                              "reason": f"font missing; render checks skipped: {', '.join(missing)}"})
                continue
            rendered_fonts = {normalized_font(font) for font in rendered.fonts}
            substituted = sorted(font for font in requested if normalized_font(font) not in rendered_fonts)
            if substituted:
                skips.append({"slide": slide_index, "shape": name,
                              "reason": "font substitution detected; render checks skipped: "
                                        f"requested={sorted(requested)}, rendered={sorted(rendered.fonts)}"})
                continue
            inspected += 1
            shape_left, shape_top = shape.left / 12700, shape.top / 12700
            shape_right = (shape.left + shape.width) / 12700
            shape_bottom = (shape.top + shape.height) / 12700
            exceeded = []
            if rendered.left < shape_left - tolerance_pt:
                exceeded.append("left")
            if rendered.top < shape_top - tolerance_pt:
                exceeded.append("top")
            if rendered.right > shape_right + tolerance_pt:
                exceeded.append("right")
            if rendered.bottom > shape_bottom + tolerance_pt:
                exceeded.append("bottom")
            if exceeded:
                issues.append(RenderIssue(
                    "render.text_overflow", slide_index, name,
                    f"bounds=({rendered.left:.1f},{rendered.top:.1f},{rendered.right:.1f},"
                    f"{rendered.bottom:.1f})pt, shape=({shape_left:.1f},{shape_top:.1f},"
                    f"{shape_right:.1f},{shape_bottom:.1f})pt, exceeded={','.join(exceeded)}",
                ))
            if rendered.bottom > float(qa["text_max_ymax_pt"]):
                issues.append(RenderIssue(
                    "render.page_text_ymax", slide_index, name,
                    f"text ymax={rendered.bottom:.1f}pt > {qa['text_max_ymax_pt']}pt",
                ))
    issues.sort(key=lambda item: (item.slide, item.rule, item.shape, item.evidence))
    return issues, skips, inspected


def run_libreoffice(path: Path, rules: dict) -> dict:
    soffice = soffice_executable()
    if soffice is None:
        return result(path, "SKIP", skips=[{"reason": "LibreOffice headless unavailable"}])
    installed = installed_non_windows_fonts()
    presentation = Presentation(path)
    with tempfile.TemporaryDirectory(prefix="render-check-") as directory:
        svg = libreoffice_svg(path, soffice, Path(directory))
        rendered = parse_svg_text(
            svg,
            presentation.slide_width / rules["units"]["emu_per_inch"],
            presentation.slide_height / rules["units"]["emu_per_inch"],
        )
    issues, skips, inspected = inspect_libreoffice_presentation(presentation, rendered, rules, installed)
    if issues:
        status = "FAIL"
    elif inspected:
        status = "PASS"
    else:
        status = "SKIP"
        skips.append({"reason": "no text shape had a verifiable, non-substituted font"})
    return result(path, status, issues, skips)


def inspect_presentation(presentation, rules, missing_heading: bool):
    qa = rules["qa"]
    tolerance_pt = rules["qa"]["canvas_overflow_tolerance_in"] * rules["units"]["pt_per_inch"]
    issues, skips = [], []
    for slide_index in range(1, presentation.Slides.Count + 1):
        slide = presentation.Slides.Item(slide_index)
        for shape in iter_shapes(slide.Shapes):
            try:
                if not bool(shape.HasTextFrame) or not bool(shape.TextFrame2.HasText):
                    continue
                frame = shape.TextFrame2
                text_range = frame.TextRange
                value = str(text_range.Text)
                if not value.strip():
                    continue
            except Exception:
                continue

            name = shape_name(shape)
            fonts = text_font_names(text_range)
            if missing_heading and rules["fonts"]["heading"] in fonts:
                skips.append({"slide": slide_index, "shape": name,
                              "reason": f"heading font missing: {rules['fonts']['heading']}"})
                continue

            left = float(text_range.BoundLeft)
            top = float(text_range.BoundTop)
            width = float(text_range.BoundWidth)
            height = float(text_range.BoundHeight)
            right, bottom = left + width, top + height
            shape_left, shape_top = float(shape.Left), float(shape.Top)
            shape_right = shape_left + float(shape.Width)
            shape_bottom = shape_top + float(shape.Height)

            exceeded = []
            if left < shape_left - tolerance_pt:
                exceeded.append("left")
            if top < shape_top - tolerance_pt:
                exceeded.append("top")
            if right > shape_right + tolerance_pt:
                exceeded.append("right")
            if bottom > shape_bottom + tolerance_pt:
                exceeded.append("bottom")
            if exceeded:
                issues.append(RenderIssue(
                    "render.text_overflow", slide_index, name,
                    f"bounds=({left:.1f},{top:.1f},{right:.1f},{bottom:.1f})pt, "
                    f"shape=({shape_left:.1f},{shape_top:.1f},{shape_right:.1f},{shape_bottom:.1f})pt, "
                    f"exceeded={','.join(exceeded)}",
                ))
            if bottom > float(qa["text_max_ymax_pt"]):
                issues.append(RenderIssue(
                    "render.page_text_ymax", slide_index, name,
                    f"text ymax={bottom:.1f}pt > {qa['text_max_ymax_pt']}pt",
                ))

    issues.sort(key=lambda item: (item.slide, item.rule, item.shape, item.evidence))
    return issues, skips


def run(path: Path, rules: dict, style: str | None = None,
        manifest_path: Path | None = None) -> dict:
    rules = style_rules(rules, style, manifest_path)
    if platform.system() != "Windows":
        try:
            return run_libreoffice(path, rules)
        except Exception as error:
            return result(path, "ERROR", error=f"{type(error).__name__}: {error}")
    try:
        import pythoncom
        import win32com.client
        import winreg
    except ImportError as error:
        return result(path, "SKIP", skips=[{"reason": f"pywin32 unavailable: {error}"}])

    installed = installed_font_names(winreg)
    body_font = rules["fonts"]["body"]
    heading_font = rules["fonts"]["heading"]
    if not font_available(body_font, installed):
        return result(path, "SKIP", skips=[{"reason": f"body font missing: {body_font}"}])
    missing_heading = not font_available(heading_font, installed)

    app = presentation = None
    pythoncom.CoInitialize()
    try:
        app = win32com.client.DispatchEx("PowerPoint.Application")
        presentation = app.Presentations.Open(str(path.resolve()), True, False, False)
        issues, skips = inspect_presentation(presentation, rules, missing_heading)
        return result(path, "FAIL" if issues else "PASS", issues, skips)
    except Exception as error:
        return result(path, "ERROR", error=f"{type(error).__name__}: {error}")
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()


def main(argv=None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--rules", type=Path, default=Path(__file__).with_name("house-rules.yaml"))
    parser.add_argument("--manifest", type=Path,
                        help="manifest used to select the render style")
    parser.add_argument("--style", help="render style when no manifest is available")
    parser.add_argument("--json", action="store_true")
    args = parser.parse_args(argv)
    try:
        if not args.pptx.is_file():
            raise FileNotFoundError(args.pptx)
        payload = run(args.pptx, load_rules(args.rules), args.style, args.manifest)
    except Exception as error:
        payload = result(args.pptx, "ERROR", error=f"{type(error).__name__}: {error}")

    if args.json:
        print(json.dumps(payload, ensure_ascii=False, indent=2))
    else:
        print(f"{payload['status']} {payload['file']}")
        for issue in payload["issues"]:
            print(f"  p{issue['slide']} {issue['rule']}: {issue['evidence']}")
        for skip in payload["skips"]:
            where = f"p{skip['slide']} {skip['shape']}: " if "slide" in skip else ""
            print(f"  SKIP {where}{skip['reason']}")
        if payload.get("error"):
            print(f"  render.error: {payload['error']}")
    return {"PASS": 0, "FAIL": 1, "ERROR": 2, "SKIP": 3}[payload["status"]]


if __name__ == "__main__":
    raise SystemExit(main())
